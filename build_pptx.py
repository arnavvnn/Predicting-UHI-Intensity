#!/usr/bin/env python3
"""
MLPR UHI Bengaluru — Final Presentation Generator
Produces a 39-slide PPTX: navy/teal/gold palette, Calibri, warm off-white body slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x15, 0x23, 0x3E)
TEAL   = RGBColor(0x1F, 0x78, 0x8A)
GOLD   = RGBColor(0xD4, 0x93, 0x2A)
OFFWHT = RGBColor(0xF6, 0xF5, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x18, 0x18, 0x18)
MID    = RGBColor(0x50, 0x50, 0x50)
LTEAL  = RGBColor(0xC8, 0xE4, 0xEA)
LGRAY  = RGBColor(0xCE, 0xCD, 0xC9)
HTEAL  = RGBColor(0xE6, 0xF4, 0xF7)
NAVYLT = RGBColor(0x1C, 0x30, 0x52)

W, H = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Low-level primitives ──────────────────────────────────────────────────────

def _blank(prs):
    for lay in prs.slide_layouts:
        if lay.name == 'Blank':
            return lay
    return prs.slide_layouts[6]

BLANK = None  # filled after prs created


def rect(slide, l, t, w, h, fill=None, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
    else:
        s.line.fill.background()
    return s


def txb(slide, text, l, t, w, h, sz=16, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.name = "Calibri"; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return b


def _run(p, text, sz, bold=False, color=DARK, italic=False):
    r = p.add_run()
    r.text = text; r.font.name = "Calibri"; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color


def bullets(slide, items, l, t, w, h, sz=15):
    """
    Multi-paragraph text box.
    '# text'   → teal bold header  (sz+1)
    '## text'  → navy bold sub-head (sz)
    '- text'   → indented sub-bullet (sz-1, MID)
    'text'     → normal bullet • (sz, DARK)
    """
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if item.startswith('# '):
            _run(p, item[2:], sz + 1, bold=True, color=TEAL)
            p.space_before = Pt(8)
        elif item.startswith('## '):
            _run(p, item[3:], sz, bold=True, color=NAVY)
            p.space_before = Pt(5)
        elif item.startswith('- '):
            _run(p, "    ‣  " + item[2:], sz - 1, color=MID)
        else:
            _run(p, "•  " + item, sz, color=DARK)
        p.space_after = Pt(2)
    return b


# ── Slide chrome helpers ──────────────────────────────────────────────────────

def _chrome(slide):
    rect(slide, 0, 0, W, H, OFFWHT)
    rect(slide, 0, 0, W, Inches(1.0), NAVY)
    rect(slide, 0, Inches(1.0), Inches(0.07), H - Inches(1.0), TEAL)
    rect(slide, 0, H - Inches(0.045), W, Inches(0.045), GOLD)


def _header(slide, text):
    txb(slide, text, Inches(0.30), Inches(0.12), Inches(12.73), Inches(0.82),
        sz=26, bold=True, color=WHITE)


# ── Slide type functions ──────────────────────────────────────────────────────

def mk_title(title1, title2, authors, course):
    sl = prs.slides.add_slide(_blank(prs))
    rect(sl, 0, 0, W, H, NAVY)
    rect(sl, 0, 0, W, Inches(0.07), TEAL)
    rect(sl, 0, H - Inches(0.07), W, Inches(0.07), GOLD)
    rect(sl, 0, Inches(2.05), W, Inches(3.4), NAVYLT)
    txb(sl, title1, Inches(0.65), Inches(2.18), Inches(12.03), Inches(1.0),
        sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, title2, Inches(0.65), Inches(3.18), Inches(12.03), Inches(0.75),
        sz=18, color=LTEAL, align=PP_ALIGN.CENTER)
    rect(sl, Inches(4.4), Inches(4.08), Inches(4.53), Inches(0.045), GOLD)
    txb(sl, authors, Inches(0.65), Inches(4.28), Inches(12.03), Inches(0.52),
        sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, course, Inches(0.65), Inches(4.86), Inches(12.03), Inches(0.42),
        sz=13, color=RGBColor(0x88, 0xB2, 0xC6), align=PP_ALIGN.CENTER)


def mk_section(num, title, subtitle=None):
    sl = prs.slides.add_slide(_blank(prs))
    rect(sl, 0, 0, W, H, NAVY)
    rect(sl, 0, 0, W, Inches(0.07), TEAL)
    rect(sl, 0, H - Inches(0.07), W, Inches(0.07), GOLD)
    rect(sl, Inches(5.65), Inches(1.62), Inches(2.03), Inches(0.65), TEAL)
    txb(sl, f"SECTION  {num}", Inches(5.65), Inches(1.64), Inches(2.03), Inches(0.61),
        sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, title, Inches(0.75), Inches(2.55), Inches(11.83), Inches(1.4),
        sz=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        txb(sl, subtitle, Inches(1.5), Inches(4.1), Inches(10.33), Inches(0.7),
            sz=17, color=LTEAL, align=PP_ALIGN.CENTER)


def mk_content(title, items, sz=15):
    sl = prs.slides.add_slide(_blank(prs))
    _chrome(sl); _header(sl, title)
    bullets(sl, items, Inches(0.30), Inches(1.08), Inches(12.73), Inches(6.2), sz)


def mk_two_col(title, left, right, lhdr=None, rhdr=None, sz=14):
    sl = prs.slides.add_slide(_blank(prs))
    _chrome(sl); _header(sl, title)
    MX = Inches(6.82)
    rect(sl, MX, Inches(1.0), Inches(0.025), H - Inches(1.0), LGRAY)
    y0 = Inches(1.08)
    if lhdr:
        txb(sl, lhdr, Inches(0.30), y0, Inches(6.3), Inches(0.42),
            sz=sz+1, bold=True, color=TEAL)
        yl = Inches(1.54)
    else:
        yl = y0
    if rhdr:
        txb(sl, rhdr, MX + Inches(0.18), y0, Inches(6.13), Inches(0.42),
            sz=sz+1, bold=True, color=TEAL)
        yr = Inches(1.54)
    else:
        yr = y0
    bullets(sl, left,  Inches(0.30), yl, Inches(6.32), Inches(5.8), sz)
    bullets(sl, right, MX + Inches(0.18), yr, Inches(6.15), Inches(5.8), sz)


def mk_table(title, headers, rows, note=None, cws=None):
    sl = prs.slides.add_slide(_blank(prs))
    _chrome(sl); _header(sl, title)
    TL, TT, TW = Inches(0.55), Inches(1.1), Inches(12.23)
    nc = len(headers)
    if cws:
        col_w = [Inches(c) for c in cws]
    else:
        cw = TW / nc
        col_w = [cw] * nc
    HH, RH = Inches(0.62), Inches(0.54)
    # header row
    x = TL
    for ci, h in enumerate(headers):
        rect(sl, x, TT, col_w[ci], HH, NAVY)
        txb(sl, h, x + Inches(0.05), TT + Inches(0.07),
            col_w[ci] - Inches(0.1), HH - Inches(0.1),
            sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += col_w[ci]
    # data rows
    for ri, row in enumerate(rows):
        bg = HTEAL if ri % 2 == 0 else WHITE
        y = TT + HH + ri * RH
        x = TL
        for ci, cell in enumerate(row):
            rect(sl, x, y, col_w[ci], RH, bg)
            txb(sl, str(cell), x + Inches(0.05), y + Inches(0.06),
                col_w[ci] - Inches(0.1), RH - Inches(0.1),
                sz=12, bold=(ci == 0), color=(NAVY if ci == 0 else DARK),
                align=PP_ALIGN.CENTER)
            x += col_w[ci]
    if note:
        yn = TT + HH + len(rows) * RH + Inches(0.1)
        txb(sl, note, TL, yn, TW, Inches(0.45), sz=12, color=MID, italic=True)


def mk_thanks(team, course):
    sl = prs.slides.add_slide(_blank(prs))
    rect(sl, 0, 0, W, H, NAVY)
    rect(sl, 0, 0, W, Inches(0.07), TEAL)
    rect(sl, 0, H - Inches(0.07), W, Inches(0.07), GOLD)
    txb(sl, "Thank You", Inches(0.75), Inches(1.75), Inches(11.83), Inches(1.65),
        sz=62, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, "Questions & Discussion",
        Inches(0.75), Inches(3.4), Inches(11.83), Inches(0.75),
        sz=22, color=LTEAL, align=PP_ALIGN.CENTER)
    rect(sl, Inches(4.2), Inches(4.28), Inches(4.93), Inches(0.045), GOLD)
    txb(sl, team, Inches(0.75), Inches(4.48), Inches(11.83), Inches(0.52),
        sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(sl, course, Inches(0.75), Inches(5.07), Inches(11.83), Inches(0.42),
        sz=13, color=RGBColor(0x88, 0xB2, 0xC6), align=PP_ALIGN.CENTER)


# =============================================================================
# BUILD THE 39 SLIDES
# =============================================================================

# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────
mk_title(
    "Predicting Urban Heat Island Severity in Bengaluru",
    "Ensemble Machine Learning on Satellite, Urban Morphology & Socioeconomic Data",
    "Arnav Nathani  ·  Sahil Aleem  ·  Vayun Gupta",
    "MLPR End-Term Project  ·  Plaksha University  ·  Spring 2026"
)

# ── SLIDE 2 — Section 1 Header ────────────────────────────────────────────────
mk_section("01", "Introduction & Motivation", "The Urban Heat Island Problem")

# ── SLIDE 3 — Why Cities Are Warming ─────────────────────────────────────────
mk_content("Why Cities Are Getting Hotter", [
    "Global average temperature has risen ~1.1 °C — but cities run 2–5 °C hotter than surroundings",
    "Urban Heat Island (UHI) effect: a persistent phenomenon in every major city on Earth",
    "# Key physical drivers",
    "Impervious surfaces (concrete, asphalt) absorb solar energy and re-radiate as heat",
    "Reduced vegetation eliminates the cooling effect of evapotranspiration",
    "Urban canyon geometry traps long-wave radiation between buildings",
    "Anthropogenic heat: vehicle exhaust, AC systems, industrial processes add direct heat",
    "# Why this matters now",
    "India's urban population will reach 600 million by 2036 (UN World Urbanisation Prospects)",
    "Heat stress responsible for 12,000+ annual deaths in India (MoHFW 2022)",
    "Every 1 °C rise in urban temperature increases electricity demand by 2–4%",
])

# ── SLIDE 4 — What is UHI? ────────────────────────────────────────────────────
mk_content("What is the Urban Heat Island Effect?", [
    "UHI = urban areas are persistently warmer than surrounding rural/suburban areas",
    "Measured via Land Surface Temperature (LST) from satellite thermal imagery",
    "# Severity classification — Liu & Zhang (2011) Mean–SD method",
    "Uses the city's own LST distribution: class boundaries at μ − 0.5σ and μ + 0.5σ",
    "- Cool Island:   LST < μ − 0.5σ       →  < 32.42 °C   →  526 pts  (21.6%)",
    "- Neutral:       μ − 0.5σ to μ + 0.5σ  →  32.42–35.14 °C  →  1,173 pts (48.2%)",
    "- Hot UHI:       LST > μ + 0.5σ       →  > 35.14 °C   →  734 pts  (30.2%)",
    "# Why raw LST (not UHII = LST_urban − LST_rural)?",
    "Our 10×10 km ROI is entirely urban — no rural reference pixel exists",
    "Raw LST is consistent with Mansouri (2025), Hoang (2025), Manna (2026)",
    "Liu & Zhang classification is defined on raw LST — same methodology as lit.",
])

# ── SLIDE 5 — Why Bengaluru? ──────────────────────────────────────────────────
mk_content("Why Bengaluru?", [
    "India's 3rd largest city: 12 M+ population, India's Silicon Valley",
    "88% loss of green cover between 1973 and 2017 (IISc remote sensing study)",
    "Rapid concretisation: wetlands converted to IT campuses, ring roads, housing",
    "Tropical savanna climate (Köppen Aw): peak heat March–May before monsoon",
    "# The research gap",
    "Conducted a survey of 11 peer-reviewed UHI ML papers (2023–2026)",
    "Zero prior papers apply ML to UHI prediction specifically for Bengaluru",
    "Ours is the first ML-based UHI study for Bengaluru in published literature",
    "# Why it matters beyond Bengaluru",
    "Bengaluru's extreme urban growth rate = extreme projected UHI risk trajectory",
    "Local government mitigation planning requires spatially resolved prediction models",
    "Provides a reproducible template for other Indian tier-1 cities",
])

# ── SLIDE 6 — Research Gaps ───────────────────────────────────────────────────
mk_content("Gaps in Existing Literature", [
    "# Validation gaps",
    "Most ML-UHI studies: single 70/30 or 80/20 random split — inflated by spatial autocorrelation",
    "Only 1 of 11 surveyed papers (Mansouri 2025) uses stratified cross-validation",
    "No paper tests geographic robustness via Spatial Block CV",
    "# Feature gaps",
    "No paper combines: socioeconomic wealth index + 3D urban morphology + remote sensing + met.",
    "No paper uses Meta Relative Wealth Index (RWI) for UHI prediction",
    "No paper derives Sky View Factor proxy + Height Variability from OSMnx",
    "# Cross-city gap",
    "No Indian-context UHI paper with empirical cross-city domain adaptation",
    "No paper quantifies the scale-mismatch failure mode in cross-city transfer",
    "# Geography gap",
    "All 11 surveyed papers: USA, Vietnam, Africa, Mumbai, Indonesia — none cover Bengaluru",
])

# ── SLIDE 7 — Our Approach ────────────────────────────────────────────────────
mk_content("Our Approach — Four Contributions", [
    "# 1. Custom dataset",
    "2,433-point dataset: 27 columns, 10×10 km ROI in central Bengaluru, 50 m grid",
    "7 heterogeneous data sources merged (GEE, ERA5, Sentinel-5P, MODIS, OSMnx, WorldPop, Meta RWI)",
    "# 2. Tuned ensemble model",
    "XGBoost + LightGBM + CatBoost, unweighted average — Optuna TPE 60 trials per model",
    "# 3. Honest evaluation",
    "Primary: Stratified 5-fold CV on LST tertiles (balanced folds)",
    "Secondary: Spatial Block CV (KMeans checkerboard, geographic independence)",
    "# 4. Cross-city validation",
    "BLR ↔ Pune domain adaptation: identified scale-mismatch problem, demonstrated fix",
    "First cross-city domain adaptation experiment for Indian UHI context",
])

# ── SLIDE 8 — Roadmap ─────────────────────────────────────────────────────────
mk_content("Project Roadmap", [
    "# 01 · Introduction & Motivation",
    "UHI problem, why Bengaluru, gaps in existing literature",
    "# 02 · Dataset & Feature Engineering",
    "Data collection pipeline, 7 sources, 27 → 20 features, Pune cross-city dataset",
    "# 03 · Methodology",
    "Model progression, ensemble architecture, Optuna tuning, CV journey",
    "# 04 · Results & Evaluation",
    "Headline metrics (R²=0.7940), residual analysis, cross-paper benchmark,",
    "- SHAP attribution, novelty ablation, cross-city transfer",
    "# 05 · Conclusions",
    "Key findings, novelty claims, limitations, future work",
], sz=16)

# ── SLIDE 9 — Section 2 Header ────────────────────────────────────────────────
mk_section("02", "Dataset & Feature Engineering",
           "From Raw Satellite Data to an ML-Ready Dataset")

# ── SLIDE 10 — Study Area ─────────────────────────────────────────────────────
mk_content("Study Area — Bengaluru ROI", [
    "10 × 10 km Region of Interest (ROI) in central Bengaluru",
    "Coordinates: 77.52–77.62 °E,  12.94–13.04 °N",
    "2,433 sample points on a regular 50 m grid",
    "Coverage: CBD, residential areas, parks (Cubbon Park, Lalbagh), industrial zones, water bodies",
    "# Why this extent?",
    "Captures the full urban gradient — from dense concrete core to greener urban fringe",
    "50 m resolution balances spatial granularity with Landsat pixel size (30 m native)",
    "All 7 data sources have ≥ 50 m native resolution → no upsampling artefacts",
    "# Cross-city: Pune",
    "Equivalent ROI in Pune: 2,382 points, same 50 m grid methodology",
    "Used exclusively for cross-city domain adaptation experiments",
])

# ── SLIDE 11 — Data Collection ────────────────────────────────────────────────
mk_content("Data Collection — 7 Sources", [
    "# Remote sensing (Google Earth Engine)",
    "Landsat 8/9 Band 10 → LST (Jiménez-Muñoz algorithm)  +  NDVI, NDBI, EVI, Albedo",
    "MODIS Terra MOD11A1 → independent LST cross-check (used as feature: MODIS_LST)",
    "# Climate reanalysis",
    "ERA5 hourly → Air_Temp_C, Relative_Humidity, Wind_Speed, Surface_Pressure_hPa, Soil_Moisture",
    "# Air quality",
    "Sentinel-5P TROPOMI → NO2_Emissions (tropospheric column density)",
    "# Urban morphology (OSMnx)",
    "Building footprints → Building_Density_Ratio, Height_Variability, SVF_Proxy (Sky View Factor)",
    "Road network → Street_Density, Dist_Highway  |  Green/water buffers → Dist_Park, Dist_Water",
    "# Population & socioeconomic",
    "WorldPop 100 m → Pop_Density  |  Meta RWI → Relative_Wealth_Index, Dist_to_RWI_Node",
])

# ── SLIDE 12 — Feature Engineering ───────────────────────────────────────────
mk_content("Feature Engineering — 27 → 20 Features", [
    "Raw dataset: 27 columns × 2,433 points",
    "# Columns dropped (7 total)",
    "LST                → regression target (cannot be a feature)",
    "Latitude, Longitude → per mentor guidance: prevents model memorising geography",
    "EVI, SAVI, MNDWI   → multicollinear with NDVI / NDBI  (Pearson |r| > 0.85)",
    "Volume_Density     → 84.8% missing values — the only null column in the dataset",
    "# Result: 20 clean features",
    "All remaining 20 columns: zero null values after dropping Volume_Density",
    "KNN imputation (k=5) applied before modelling as safety net for residual nulls",
    "StandardScaler for linear/distance models; raw features for tree ensembles (XGB/LGB/CB)",
])

# ── SLIDE 13 — Feature Categories Table ──────────────────────────────────────
mk_table("The 20 Features — Category Breakdown",
    ["Category", "Features"],
    [
        ["Remote Sensing",          "NDVI,  NDBI"],
        ["Meteorological",          "Air_Temp_C,  Relative_Humidity,  Wind_Speed,  Soil_Moisture,  Surface_Pressure_hPa,  MODIS_LST"],
        ["Anthropogenic",           "NO2_Emissions,  Pop_Density"],
        ["OSM 3D Morphology",       "Building_Density_Ratio,  Height_Variability,  SVF_Proxy"],
        ["OSM Proximity",           "Dist_Water,  Dist_Park,  Dist_Highway,  Street_Density"],
        ["Socioeconomic  ★ Novel",  "Relative_Wealth_Index,  Dist_to_RWI_Node"],
        ["Radiative",               "Albedo"],
    ],
    note="★ = features not used in any prior UHI ML paper in our 11-paper survey",
    cws=[3.4, 8.83]
)

# ── SLIDE 14 — Target Variable ────────────────────────────────────────────────
mk_content("Target Variable & UHI Classification Scheme", [
    "# Regression target: raw Land Surface Temperature (LST, °C)",
    "Landsat Band 10 thermal retrieval — Jiménez-Muñoz single-channel algorithm",
    "BLR LST range: 27.7 – 40.6 °C  |  Mean = 33.78 °C  |  σ = 2.72 °C",
    "# Why raw LST, not UHII = LST_urban − LST_rural?",
    "Our ROI is entirely urban — no rural reference pixel available for UHII computation",
    "All comparison papers (Mansouri, Hoang, Manna) also regress on raw LST",
    "Liu & Zhang (2011) classification is defined on raw LST distribution",
    "# UHI severity classes — Liu & Zhang (2011) Mean–SD method",
    "Boundaries at μ − 0.5σ  and  μ + 0.5σ  of the city's own LST distribution",
    "- Cool Island:  LST < 32.42 °C          →  526 pts  (21.6%)",
    "- Neutral:      32.42 – 35.14 °C        →  1,173 pts (48.2%)",
    "- Hot UHI:      LST > 35.14 °C          →  734 pts  (30.2%)",
])

# ── SLIDE 15 — Pune Dataset ───────────────────────────────────────────────────
mk_content("Cross-City Dataset — Pune", [
    "2,382 points, 36 columns (same core + LULC one-hot columns + pre-computed UHII)",
    "Same collection pipeline: GEE, ERA5, OSMnx, WorldPop, Meta RWI",
    "# Critical difference from Bengaluru",
    "Pune LST is already z-scored → mean ≈ 0, std ≈ 1  (unit variance)",
    "Bengaluru LST is raw °C → mean = 33.78, std = 2.72",
    "This scale mismatch is the root cause of naive transfer failing (R² = −0.24)",
    "# All 20 BLR features are present in Pune (verified — 0 nulls)",
    "NO2_Emissions, Height_Variability, Surface_Pressure_hPa all confirmed in Pune CSV",
    "Cross-city experiments use the full 20-feature shared set (not a reduced subset)",
    "NDBI, NDVI, NO2, Pop_Density, RH, Wind_Speed, Building_Density_Ratio,",
    "- Dist_Water, Dist_Park, Dist_Highway, Street_Density, Height_Var, SVF_Proxy,",
    "- Air_Temp_C, Albedo, MODIS_LST, Soil_Moisture, Surface_Pressure, RWI, Dist_to_RWI_Node",
])

# ── SLIDE 16 — Section 3 Header ───────────────────────────────────────────────
mk_section("03", "Methodology", "From Dataset to Validated Ensemble Model")

# ── SLIDE 17 — ML Pipeline Overview ──────────────────────────────────────────
mk_content("ML Pipeline — End to End", [
    "# Step 1 · Preprocessing",
    "KNN imputation (k=5) for residual missing values (safety net)",
    "StandardScaler for linear / distance-based models only  (SVM, KNN, Linear, MLP)",
    "Tree ensembles trained on raw features — no scaling needed",
    "# Step 2 · Model selection sweep",
    "8 model families: Linear, KNN, SVM, Random Forest, XGBoost, LightGBM, CatBoost, MLP",
    "MLP architecture: Input → 64 ReLU → 32 ReLU → Output, Dropout 0.2, Adam, early stopping",
    "# Step 3 · Ensemble construction",
    "Final model: unweighted average of XGBoost + LightGBM + CatBoost predictions",
    "# Step 4 · Evaluation",
    "Primary: Stratified 5-fold CV on LST tertile bins",
    "Secondary: Spatial Block CV (KMeans checkerboard, 10 geographic folds)",
])

# ── SLIDE 18 — Model Progression ─────────────────────────────────────────────
mk_content("Model Progression — Baseline to Headline Ensemble", [
    "# Baseline models  (verified, Stratified 5-fold CV)",
    "Linear Regression:          R² = 0.4241  (severe underfitting — LST is strongly nonlinear)",
    "KNN (k=7, distance-weighted): R² = 0.5985  (overfits to local spatial clusters)",
    "MLP (64→32 ReLU):           R² = 0.5066  (neural net struggles with small n=2,433)",
    "Random Forest (n=500):      R² = 0.6943  (strong baseline, low variance)",
    "# Tree ensemble models  (tuned, verified)",
    "XGBoost  (Optuna-tuned):    R² = 0.7869  (boosting + regularisation wins)",
    "LightGBM (Optuna-tuned):    R² = 0.7855",
    "CatBoost (Optuna-tuned):    R² = 0.7875",
    "# Ensemble",
    "XGB + LGB + CatBoost (Optuna-tuned):    R² = 0.7940  ← headline",
    "Note: untuned single-model and pairwise ensemble R² values from notebook exploration",
])

# ── SLIDE 19 — Ensemble Architecture ─────────────────────────────────────────
mk_content("Ensemble Architecture — Diversity by Design", [
    "Three gradient-boosted tree models in an unweighted average ensemble",
    "# Why these three?  Diverse inductive biases",
    "XGBoost:   symmetric depth-wise trees, L1/L2 regularisation, column subsampling",
    "LightGBM:  leaf-wise growth (faster, handles larger n), histogram binning",
    "CatBoost:  ordered boosting (prevents target leakage), strong on numeric skew",
    "# Why unweighted average?",
    "Optimised stacking adds a meta-learner — overfits on 2,433 points",
    "Equal weighting performs on par with stacking here; simpler and reproducible",
    "# Ensemble gain quantified  (verified)",
    "All three models tightly clustered: XGB=0.7869, LGB=0.7855, CB=0.7875",
    "Ensemble:  R² = 0.7940  →  +0.0065 over best single model (CatBoost)",
    "Primary ensemble benefit: lower RMSE (1.233 vs ~1.252) and more stable predictions",
])

# ── SLIDE 20 — Hyperparameter Tuning ─────────────────────────────────────────
mk_content("Hyperparameter Tuning — Optuna TPE", [
    "Algorithm: Tree-structured Parzen Estimator (TPE) — Bayesian optimisation",
    "60 trials per model; objective = mean 5-fold CV R²  (not held-out test — avoids leakage)",
    "Hyperparameters locked in verify_experiments.py for full reproducibility",
    "# Locked best parameters",
    "XGB:   n_est=929,  max_depth=7,   lr=0.0194, subsample=0.870, colsample=0.632,",
    "-       min_child_weight=6,  α=1.4×10⁻⁵,  λ=2×10⁻⁷",
    "LGB:   n_est=1244, num_leaves=51, lr=0.0318, subsample=0.870, colsample=0.620,",
    "-       min_child_samples=11, α=3.3×10⁻⁶, λ=3×10⁻⁸",
    "CB:    iterations=1304, depth=7,  lr=0.0735, l2=4.285, subsample=0.877",
    "# Key convergence pattern",
    "All three models → similar subsample ≈ 0.87, colsample ≈ 0.62–0.63",
    "High n_estimators (929–1304) + low lr = slow, careful learning (reduced overfitting)",
])

# ── SLIDE 21 — Validation Strategy ───────────────────────────────────────────
mk_content("Validation Strategy — Why Rigor Matters", [
    "# Stratified 5-fold CV  (headline protocol)",
    "Folds balanced by LST tertile bins → equal Cool / Neutral / Hot in every fold",
    "Prevents class-imbalance artefacts; 5 independent estimates → stable mean ± std",
    "# Spatial Block CV  (geographic robustness check)",
    "KMeans (k=10) clusters points into 10 geographic blocks",
    "Train on 9 blocks, test on 1 — guarantees spatial independence between train/test",
    "Result: R² ≈ 0.72–0.74  (small drop; model generalises beyond training geography)",
    "# What we rejected and why",
    "Simple 80/20 random split: inflated by spatial autocorrelation → +0.03–0.05 R²",
    "Leave-One-Out CV: computationally prohibitive, very high variance on n=2,433",
    "# Cross-city validation",
    "Train BLR, test Pune (and vice versa): hardest generalisation test",
])

# ── SLIDE 22 — CV Journey ─────────────────────────────────────────────────────
mk_content("The CV Journey — How Validation Evolved", [
    "# Stage 1 — Random 80/20 split  (rejected)",
    "First approach: fast, seemingly standard",
    "Problem: spatial autocorrelation → test points near train points → inflated R²",
    "Estimated overestimation: +0.03–0.05 R² vs honest estimate",
    "# Stage 2 — Spatial Block CV  (geographic robustness)",
    "KMeans geographic blocks enforce spatial independence",
    "First honest estimate: R² ≈ 0.72–0.74; confirmed model learns physics, not coordinates",
    "# Stage 3 — Stratified 5-fold CV  ← headline",
    "Tertile-balanced folds: each fold has same Cool / Neutral / Hot proportions",
    "5× more stable estimates than single split; matches published standard (Mansouri 2025)",
    "# Stage 4 — Cross-city transfer  (hardest test)",
    "Train BLR, test Pune — different climate, scale, and preprocessing",
    "Lesson: each validation upgrade costs ~0.02–0.04 R² but earns credibility",
])

# ── SLIDE 23 — Cross-City Pipeline ───────────────────────────────────────────
mk_content("Cross-City Transfer — Design & Pipeline", [
    "# The scale-mismatch problem",
    "Bengaluru LST: raw °C  →  mean = 33.78,  σ = 2.72",
    "Pune LST: pre-standardised z-score  →  mean ≈ 0,  σ ≈ 1",
    "Naive transfer (raw features, z-scored targets):  R² = −0.2443  ← feature scale mismatch",
    "# Fix: per-city z-score standardisation of features",
    "Standardise each city's features independently (both targets already z-scored)",
    "After fix — directional BLR → Pune transfer:  R² = +0.1662  (2.4× better than old approach)",
    "# Why does z-scoring features help so much?",
    "BLR features are raw physical units; Pune features are pre-standardised",
    "Per-city z-score aligns the feature distributions before the model sees them",
    "# Best approach: combined model",
    "Pool BLR + Pune together after per-city z-score of features + targets",
    "Combined 5-fold CV:  R² = 0.7497  (vs old 18-feat result of 0.7005)",
])

# ── SLIDE 24 — Section 4 Header ───────────────────────────────────────────────
mk_section("04", "Results & Evaluation",
           "Verified Metrics from Stratified 5-Fold CV")

# ── SLIDE 25 — Headline Regression ───────────────────────────────────────────
mk_content("Headline Regression Results", [
    "# Primary metric: Stratified 5-fold CV  ·  2,433 points  ·  20 features",
    "Ensemble (XGB + LGB + CB, Optuna-tuned):   R² = 0.7940",
    "RMSE = 1.233 °C    |    MAE = 0.887 °C",
    "Normalised RMSE  =  RMSE / σ_LST  =  1.233 / 2.72  =  0.453",
    "# Interpretation",
    "The ensemble explains 79.4% of Land Surface Temperature variance across Bengaluru",
    "Average prediction error < 1 °C — practically actionable for urban planning",
    "# Reproducibility",
    "All numbers from verify_experiments.py — single script, ~4 min runtime on laptop",
    "Hyperparameters locked; Optuna not re-run — results deterministic (random_state=42)",
    "verification_summary.json is the machine-readable ground truth for all metrics",
])

# ── SLIDE 26 — Per-Model Comparison ──────────────────────────────────────────
mk_table("Per-Model Performance — Stratified 5-Fold CV",
    ["Model", "R²  (mean ± std)", "RMSE (°C)", "MAE (°C)"],
    [
        ["XGBoost  (tuned)",            "0.7869  ±  0.011", "1.254", "0.904"],
        ["LightGBM  (tuned)",           "0.7855  ±  0.014", "1.258", "0.911"],
        ["CatBoost  (tuned)",           "0.7875  ±  0.009", "1.252", "0.900"],
        ["Ensemble  (XGB + LGB + CB)",  "0.7940  ±  0.010", "1.233", "0.887"],
        ["Ensemble gain vs best single","+ 0.0065  (lower RMSE, more stable)", "—", "—"],
    ],
    note="Classification (3-class, same ensemble):  Accuracy = 0.7604  |  F1-weighted = 0.7599  |  ROC-AUC (OvR) = 0.8946",
    cws=[4.23, 3.5, 2.25, 2.25]
)

# ── SLIDE 27 — Residual Analysis ─────────────────────────────────────────────
mk_content("Residual Analysis", [
    "Residuals = actual LST − predicted LST  (°C)",
    "Distribution approximately normal; Shapiro-Wilk p < 0.05 (expected at n=2,433)",
    "Skewness = −0.61  (slight left skew)",
    "# What the skew tells us",
    "Model slightly overpredicts in cool zones → negative residuals → left skew",
    "Cool Island is the minority class (21.6%) — statistically harder to predict precisely",
    "# Tail behaviour",
    "Slight underprediction at very high LST (> 38 °C): hottest micro-zones are sparse in training",
    "No systematic bias in central LST range (32–37 °C) — covers 70% of all data points",
    "# Q-Q plot",
    "Residuals track the normal line well until ±2σ; heavier left tail beyond that",
    "Standard behaviour for tree ensemble models on real-world geospatial data",
])

# ── SLIDE 28 — Classification Results ────────────────────────────────────────
mk_content("3-Class UHI Classification Results", [
    "Classes: Cool Island  |  Neutral  |  Hot UHI  — Liu & Zhang (2011)",
    "# Headline metrics  (Stratified 5-fold CV)",
    "Accuracy         =  0.7604",
    "F1-weighted      =  0.7599",
    "ROC-AUC (OvR)   =  0.8946",
    "# Per-class insight",
    "Cool Island and Hot UHI: better F1 (clear LST signal at distribution tails)",
    "Neutral class: hardest — boundary bleed ±1 °C with adjacent classes",
    "Structural reason: boundaries (32.42°C, 35.14°C) are ±0.5σ → only 2.72 °C wide band",
    "# Why run both regression and classification?",
    "Regression gives quantitative LST prediction (°C) for policy and engineering use",
    "Classification matches Mansouri (2025) protocol — enables direct apples-to-apples comparison",
])

# ── SLIDE 29 — Cross-Paper Benchmark ─────────────────────────────────────────
mk_table("Cross-Paper Benchmark",
    ["Paper", "Region", "Best R²  /  Acc", "Protocol"],
    [
        ["Hoang & Nguyen 2025",     "Da Nang (VN)",       "R² = 0.90",             "Single 70/30 split"],
        ["Lynda et al. 2025",       "Africa (10 K cities)","R² = 0.84",             "10-fold CV"],
        ["Manna et al. 2026",       "Mumbai (India)",      "R² = 0.80",             "In-sample"],
        ["Kusumadewi et al. 2025",  "Malang (Indonesia)",  "R² = 0.81",             "Held-out split"],
        ["Mansouri & Erfani 2025",  "USA Midwest",         "Acc=0.76, AUC=0.91",   "Stratified CV"],
        ["Ours  ★",                 "Bengaluru",           "R²=0.7940, AUC=0.8946","Stratified 5-fold ✓"],
    ],
    note="★ Norm-RMSE = RMSE/σ_LST = 0.453 — only fair cross-paper metric (different protocols invalidate direct R² comparison)",
    cws=[3.3, 2.8, 3.0, 3.13]
)

# ── SLIDE 30 — Spatial Block CV ───────────────────────────────────────────────
mk_content("Spatial Block CV — Geographic Robustness", [
    "KMeans (k=10) clusters all 2,433 Bengaluru points into 10 geographic blocks",
    "Train on 9 blocks, test on 1 — no test point is geographically near any training point",
    "Result: R² ≈ 0.72–0.74  (vs 0.7940 stratified) — small drop, model genuinely generalises",
    "# What Spatial Block CV validates",
    "Model learns physical relationships (heat, vegetation, morphology), not coordinate interpolation",
    "Landsat grid data has strong spatial autocorrelation — independence test is critical",
    "# Why Stratified 5-fold remains the headline",
    "Spatial Block CV: only 10% of data in each test fold → higher variance estimate",
    "Stratified 5-fold: 20% test per fold, 5 estimates → published standard for reporting",
    "# Comparison to literature",
    "Only Mansouri (2025) among our 11 surveyed papers uses comparable validation rigour",
    "Most published R² values are inflated vs honest stratified estimates",
])

# ── SLIDE 31 — Within-City CV ─────────────────────────────────────────────────
mk_content("Within-City CV — BLR vs Pune on 18 Common Features", [
    "Experiment: run same 5-fold CV protocol on each city independently with 18 common features",
    "# Results",
    "BLR alone   (20 features, z-scored, 2,433 pts):  R² = 0.7934",
    "Pune alone  (20 features, z-scored, 2,382 pts):  R² = 0.7194",
    "BLR alone   (20 features, raw target):            R² = 0.7940  (z-scoring barely affects trees)",
    "# Why is Pune lower than BLR?",
    "Pune LST is z-scored → less absolute variance → tree splits carry less information",
    "Pune has a smoother spatial LST distribution — fewer extreme hotspots to learn from",
    "# Combined model (BLR + Pune together): R² = 0.7497",
    "Old result with 18 features was 0.7005 — adding the 3 missing features adds +0.049",
    "Both cities above 0.71 R² individually; combined pools 4,815 points across both cities",
    "# Takeaway: using all 20 shared features matters — the 18-feat COMMON was leaving signal on the table",
])

# ── SLIDE 32 — Cross-City Transfer Table ─────────────────────────────────────
mk_table("Cross-City Transfer — All Configurations",
    ["Configuration", "Features", "R²"],
    [
        ["Naive BLR→Pune (raw features, z-scored targets)", "20", "−0.2443"],
        ["Naive Pune→BLR (raw features, z-scored targets)", "20", "−0.1232"],
        ["Per-city z-score BLR→Pune",                       "20", "+0.1662"],
        ["Per-city z-score Pune→BLR",                       "20", "+0.1273"],
        ["BLR alone, within-city 5-fold CV (raw target)",   "20", " 0.7940"],
        ["BLR alone, within-city CV (z-scored target)",     "20", " 0.7934"],
        ["Pune alone, within-city 5-fold CV",               "20", " 0.7194"],
        ["Combined model (BLR + Pune), 5-fold CV",          "20", " 0.7497"],
    ],
    note="Key lesson: per-city z-score of features is mandatory. Using full 20-feat set adds +0.049 R² vs old 18-feat combined (0.7005→0.7497)",
    cws=[6.73, 2.5, 3.0]
)

# ── SLIDE 33 — Novelty Ablation ───────────────────────────────────────────────
mk_content("Group-wise Novelty Ablation", [
    "Question: are the novel features genuinely necessary, or just redundant with standard ones?",
    "# Protocol",
    "Remove all 11 novel features from the model",
    "- RWI, Dist_to_RWI_Node, NO2_Emissions, SVF_Proxy, Height_Variability,",
    "- Building_Density_Ratio, Street_Density, Dist_Highway, Dist_Park, Dist_Water, Albedo",
    "Retrain ensemble on remaining 9 standard features only",
    "# Result",
    "Standard features only (9 features):  R² = 0.6605",
    "Full model         (20 features):      R² = 0.7940",
    "ΔR²  =  +0.1335  →  novel features explain 17% of total LST variance (irreplaceable)",
    "# SHAP (43.4%) vs Ablation (17%) — two different questions",
    "SHAP = how much does the trained model rely on novel features  (attribution)",
    "Ablation = how much variance is irreplaceable if removed  (necessity)",
    "Gap: standard features carry some of the same signal → correlated information",
])

# ── SLIDE 34 — SHAP ───────────────────────────────────────────────────────────
mk_content("SHAP Feature Attribution — TreeSHAP", [
    "TreeSHAP: exact Shapley values for tree ensembles (no approximation needed)",
    "Metric: mean |SHAP value| = average absolute contribution to each prediction",
    "# Top 8 features by attribution  (verified from shap_results.csv)",
    "1.  NDBI                   20.53%  (built-up surfaces absorb & re-radiate heat — dominant)",
    "2.  Relative_Humidity      12.73%  (moisture suppresses LST peaks via latent heat)",
    "3.  Relative_Wealth_Index  12.28%  ★ novel — richer areas → more AC → heat rejection",
    "4.  NO2_Emissions           9.64%  ★ novel — anthropogenic combustion heat proxy",
    "5.  MODIS_LST               5.81%  (cross-sensor LST corroboration)",
    "6.  Soil_Moisture           5.39%  (latent heat flux cooling signal)",
    "7.  Dist_Water              5.32%  ★ novel — proximity to water bodies cools LST",
    "8.  Dist_Park               4.08%  ★ novel — green space proximity cooling effect",
    "# Novel features in aggregate",
    "Total SHAP attribution from all 11 novel features:  43.4%",
    "RWI alone (12.28%) outranks NDVI (3.76%) — socioeconomics is a genuine independent driver",
])

# ── SLIDE 35 — Section 5 Header ───────────────────────────────────────────────
mk_section("05", "Conclusions", "What We Found, Why It Matters, What Comes Next")

# ── SLIDE 36 — Key Findings ───────────────────────────────────────────────────
mk_content("Key Findings", [
    "# Regression  (Stratified 5-fold CV, 2,433 points, 20 features)",
    "Ensemble achieves  R² = 0.7940,  RMSE = 1.233 °C,  MAE = 0.887 °C",
    "Ensemble outperforms best single model (CB=0.7875) by +0.0065 R²; main gain is lower RMSE & stability",
    "# Novel features are genuinely necessary",
    "Group-wise ablation:  ΔR² = +0.1335  →  17% of total LST variance explained by novel feats",
    "SHAP attribution:  novel features account for 43.4% of all model attribution",
    "RWI alone is the #3 most important feature (12.28%) — socioeconomics drive UHI",
    "# Classification",
    "3-class UHI severity:  Accuracy = 0.7604,  F1 = 0.7599,  AUC = 0.8946",
    "Matches Mansouri & Erfani (2025) on Acc/F1 with 4.4× less data; AUC −0.015",
    "# Cross-city transfer",
    "Per-city z-score alignment fixes negative transfer (−0.24 → +0.17 → 0.75 combined)",
])

# ── SLIDE 37 — Novelty Claims ─────────────────────────────────────────────────
mk_content("Novelty Claims", [
    "# 1.  First ML-based UHI study for Bengaluru",
    "Zero prior papers in our 11-paper survey cover Bengaluru",
    "# 2.  Novel feature combination",
    "Only paper jointly combining Meta RWI + OSMnx 3D morphology + remote sensing + meteorological",
    "# 3.  Cross-city domain adaptation for Indian cities",
    "Only Indian-context UHI paper with empirically demonstrated cross-city transfer",
    "Identified the feature scale-mismatch failure mode; fix: per-city z-score → R² −0.24 → +0.17 → 0.75",
    "# 4.  Quantified novelty contribution",
    "Ablation: ΔR² = +0.1335 from novel features (17% of total variance — irreplaceable)",
    "SHAP: novel features = 43.4% of attribution; RWI (12.28%) outranks NDVI (3.76%)",
    "# 5.  Honest, reproducible validation",
    "Stratified 5-fold CV + Spatial Block CV — comparable rigour to only Mansouri (2025)",
    "Single verify_experiments.py script reproduces all reported numbers in ~4 minutes",
])

# ── SLIDE 38 — Limitations & Future Work ─────────────────────────────────────
mk_two_col("Limitations & Future Work",
    left=[
        "# Current limitations",
        "Single dry-season snapshot — captures summer-peak UHI only, not seasonal variation",
        "Small dataset (2,433 pts) relative to global UHI studies (5K–10K+)",
        "Volume_Density unusable: 84.8% missing — 3D building volume excluded",
        "No probabilistic uncertainty: point predictions only",
        "Directional transfer is positive but modest (+0.17 R²)",
        "Single-city training; cross-city generalisation still limited",
    ],
    right=[
        "# Future work",
        "Multi-temporal dataset: monthly Landsat composites → seasonal UHI tracking",
        "Conformal prediction: calibrated uncertainty intervals per prediction",
        "Fine-grained road features: lane count, traffic volume, surface material",
        "Expand to pan-India model: 5 cities with unified z-score pipeline",
        "Causal analysis: which features can be targeted by policy interventions?",
        "Probabilistic classification: predict UHI class probabilities, not just labels",
    ],
    sz=14
)

# ── SLIDE 39 — Thank You ──────────────────────────────────────────────────────
mk_thanks(
    "Arnav Nathani  ·  Sahil Aleem  ·  Vayun Gupta",
    "MLPR End-Term Project  ·  Plaksha University  ·  Spring 2026"
)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "/Users/arnavnathani/Desktop/MLPR/Project/dataset and training/MLPR_UHI_Bengaluru_Final_Presentation.pptx"
prs.save(OUT)
print(f"Saved  →  {OUT}")
print(f"Slides : {len(prs.slides)}")
